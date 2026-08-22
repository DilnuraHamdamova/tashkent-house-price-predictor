"""Build the defense deck from the repository's verified results."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "presentation" / "Tashkent_House_Price_Defense.pptx"

NAVY = RGBColor(5, 24, 70)
BLUE = RGBColor(13, 91, 255)
CYAN = RGBColor(43, 184, 255)
LIGHT = RGBColor(240, 246, 255)
MID = RGBColor(99, 117, 153)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(13, 157, 89)
YELLOW = RGBColor(232, 157, 0)
RED = RGBColor(208, 55, 72)


def set_run(run, size: int, *, color=NAVY, bold: bool = False) -> None:
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 20,
    color=NAVY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    set_run(run, size, color=color, bold=bold)
    return box


def add_box(slide, x: float, y: float, w: float, h: float, *, fill=WHITE, line=BLUE, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    return shape


def add_header(slide, number: int, section: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_box(slide, 0.45, 0.30, 3.35, 0.50, fill=WHITE, line=BLUE)
    add_text(slide, f"AI/ML CAPSTONE  /  {section}", 0.68, 0.39, 2.95, 0.25, size=12, bold=True)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.10), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()
    for x, y, diameter in [(11.8, -0.2, 1.8), (12.5, 6.7, 1.0), (-0.3, 6.8, 1.1)]:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(diameter), Inches(diameter)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = LIGHT
        circle.fill.transparency = 30
        circle.line.color.rgb = RGBColor(222, 231, 246)
    add_text(
        slide,
        f"Slide {number}",
        11.95,
        0.35,
        0.85,
        0.28,
        size=13,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, 0.48, 1.02, 12.3, 0.72, size=29, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.52, 1.75, 12.0, 0.46, size=15, color=MID)


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    heading: str,
    body: str,
    *,
    accent=BLUE,
    heading_size: int = 17,
    body_size: int = 14,
) -> None:
    add_box(slide, x, y, w, h, fill=WHITE, line=accent)
    marker = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.18), Inches(0.34), Inches(0.34)
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = accent
    marker.line.fill.background()
    add_text(
        slide,
        heading,
        x + 0.62,
        y + 0.18,
        w - 0.78,
        0.35,
        size=heading_size,
        bold=True,
        color=accent,
    )
    add_text(slide, body, x + 0.22, y + 0.70, w - 0.44, h - 0.84, size=body_size)


def add_metric(slide, x: float, y: float, label: str, value: str, *, color=BLUE) -> None:
    add_box(slide, x, y, 2.35, 1.18, fill=LIGHT, line=color)
    add_text(
        slide,
        value,
        x + 0.10,
        y + 0.18,
        2.15,
        0.44,
        size=25,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, label, x + 0.10, y + 0.75, 2.15, 0.32, size=10, color=MID, align=PP_ALIGN.CENTER
    )


def build() -> Path:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    # Slide 1 — opening
    slide = presentation.slides.add_slide(blank)
    add_header(slide, 1, "INDIVIDUAL PROJECT")
    add_text(slide, "TASHKENT", 0.55, 1.35, 7.2, 0.75, size=40, bold=True)
    add_text(slide, "APARTMENT PRICE", 0.55, 2.10, 7.2, 0.82, size=39, bold=True, color=BLUE)
    add_text(slide, "PREDICTOR", 0.55, 2.92, 7.2, 0.78, size=41, bold=True)
    add_text(
        slide,
        "Dilnura Hamdamova  •  AI/ML Fundamentals Capstone",
        0.60,
        3.95,
        7.3,
        0.45,
        size=17,
        color=MID,
    )
    add_box(slide, 8.10, 1.37, 4.55, 3.15, fill=LIGHT, line=BLUE)
    add_text(slide, "PROJECT IN ONE LINE", 8.48, 1.76, 3.8, 0.35, size=14, color=BLUE, bold=True)
    add_text(
        slide,
        "Estimate a current August 2026 Tashkent apartment asking price from observable listing attributes.",
        8.48,
        2.28,
        3.67,
        1.48,
        size=19,
        bold=True,
    )
    add_text(
        slide,
        "Reference estimate — not an appraisal",
        8.48,
        4.04,
        3.65,
        0.28,
        size=12,
        color=RED,
        bold=True,
    )
    add_metric(slide, 0.65, 5.35, "MODELING ROWS", "4,214")
    add_metric(slide, 3.25, 5.35, "SELECTED MODEL", "RF")
    add_metric(slide, 5.85, 5.35, "TEST R²", "0.681", color=GREEN)
    add_metric(slide, 8.45, 5.35, "TEST MAE", "$27,195", color=GREEN)

    # Slide 2 — user and ML task
    slide = presentation.slides.add_slide(blank)
    add_header(slide, 2, "PROBLEM + TASK")
    add_title(
        slide,
        "WHO NEEDS IT — AND WHAT DOES ML DO?",
        "Lock the user, input, output, and boundary in under 45 seconds.",
    )
    add_card(
        slide,
        0.55,
        2.35,
        3.85,
        2.72,
        "USER",
        "Buyers, sellers, agents, and analysts who need a consistent August 2026 reference estimate.",
    )
    add_card(
        slide,
        4.74,
        2.35,
        3.85,
        2.72,
        "SUPERVISED REGRESSION",
        "Inputs: district, size, rooms, apartment level, building levels, new-build/resale.\n\nTarget: asking price.",
        accent=CYAN,
    )
    add_card(
        slide,
        8.93,
        2.35,
        3.85,
        2.72,
        "OUTPUT + BOUNDARY",
        "One estimated August 2026 asking price in USD, plus warnings.\n\nNot a completed sale price, guarantee, or legal appraisal.",
        accent=GREEN,
    )
    add_box(slide, 1.65, 5.55, 10.0, 0.85, fill=LIGHT, line=BLUE)
    add_text(
        slide,
        "RAW APARTMENT  →  VALIDATED PIPELINE  →  USD ESTIMATE + WARNING",
        1.95,
        5.79,
        9.4,
        0.30,
        size=19,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )

    # Slide 3 — data and method
    slide = presentation.slides.add_slide(blank)
    add_header(slide, 3, "DATA + APPROACH")
    add_title(
        slide,
        "FROM 4,867 LISTINGS TO A GROUP-SAFE TEST",
        "Privacy-minimized public HATA snapshot collected 22 August 2026.",
    )
    steps = [
        ("SOURCE", "4,867 parsed\n11 districts"),
        ("CLEAN", "257 invalid + 396\nduplicates removed"),
        ("SPLIT", "3,840 groups\n80% / 20%"),
        ("SELECT", "5-fold CV\ngroup-safe"),
        ("EVALUATE", "835 unseen rows\n768 groups"),
    ]
    for index, (heading, body) in enumerate(steps):
        x = 0.52 + index * 2.55
        add_card(slide, x, 2.32, 2.15, 2.20, heading, body, heading_size=14, body_size=16)
        if index < len(steps) - 1:
            add_text(
                slide,
                "→",
                x + 2.18,
                3.08,
                0.34,
                0.35,
                size=24,
                color=BLUE,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
    add_box(slide, 0.73, 5.05, 11.85, 1.12, fill=LIGHT, line=BLUE)
    add_text(slide, "Leakage controls", 1.05, 5.32, 1.75, 0.32, size=16, color=BLUE, bold=True)
    add_text(
        slide,
        "feature+target deduplication  •  identical fingerprints grouped  •  no price/m² leakage  •  preprocessing inside CV  •  test not used for selection",
        2.75,
        5.26,
        9.35,
        0.52,
        size=15,
    )

    # Slide 4 — experiments
    slide = presentation.slides.add_slide(blank)
    add_header(slide, 4, "MODELING + EXPERIMENTS")
    add_title(
        slide,
        "FOUR APPROACHES — ONE EVIDENCE-BASED CHOICE",
        "Final model selected by the lowest five-fold development CV MAE, not test performance.",
    )
    rows, cols = 5, 4
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(0.75), Inches(2.35), Inches(8.05), Inches(3.35)
    )
    table = table_shape.table
    widths = [2.65, 1.80, 1.80, 1.80]
    for i, width in enumerate(widths):
        table.columns[i].width = Inches(width)
    values = [
        ["MODEL", "CV MAE", "CV RMSE", "CV R²"],
        ["Median baseline", "$55,604", "$123,587", "-0.068"],
        ["Log Ridge", "$38,301", "$105,755", "0.219"],
        ["Random Forest", "$31,298", "$80,394", "0.556"],
        ["Gradient Boosting", "$33,545", "$94,081", "0.396"],
    ]
    for row in range(rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text = values[row][col]
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.10)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if row == 0 else (LIGHT if row == 3 else WHITE)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT if col == 0 else PP_ALIGN.CENTER
                for run in paragraph.runs:
                    set_run(
                        run,
                        14,
                        color=WHITE if row == 0 else (BLUE if row == 3 else NAVY),
                        bold=row in (0, 3),
                    )
    add_card(
        slide,
        9.25,
        2.35,
        3.25,
        3.35,
        "WHY RANDOM FOREST?",
        "Best validation MAE\n\nCaptures non-linear size × location interactions\n\nTrade-off: larger and less interpretable than Ridge",
        accent=GREEN,
        body_size=16,
    )

    # Slide 5 — protected test and weakness
    slide = presentation.slides.add_slide(blank)
    add_header(slide, 5, "RESULT + WEAKNESS")
    add_title(
        slide,
        "UNSEEN-TEST RESULT — WITH THE FAILURE VISIBLE",
        "The baseline comparison and error case matter more than a polished success claim.",
    )
    add_metric(slide, 0.65, 2.35, "MAE", "$27,195", color=GREEN)
    add_metric(slide, 3.15, 2.35, "RMSE", "$58,887", color=GREEN)
    add_metric(slide, 5.65, 2.35, "R²", "0.681", color=GREEN)
    add_metric(slide, 8.15, 2.35, "MAPE", "24.58%", color=GREEN)
    add_box(slide, 10.65, 2.35, 2.05, 1.18, fill=LIGHT, line=BLUE)
    add_text(
        slide,
        "−46.6%",
        10.78,
        2.53,
        1.78,
        0.40,
        size=24,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, "MAE vs baseline", 10.78, 3.07, 1.78, 0.22, size=12, color=MID, align=PP_ALIGN.CENTER
    )
    add_card(
        slide,
        0.65,
        4.05,
        7.05,
        2.02,
        "LARGEST ERROR",
        "Shayhontohur • 160 m² • 3 rooms • new build\nActual: $1,000,000   Predicted: $234,461\nAbsolute error: $765,539",
        accent=RED,
        body_size=18,
    )
    add_card(
        slide,
        8.05,
        4.05,
        4.65,
        2.02,
        "WHAT IT TEACHES",
        "Luxury premiums and condition are not observed. Mirobod and Shayhontohur MAE is high; some district slice R² values are negative.",
        accent=YELLOW,
        body_size=16,
    )

    # Slide 6 — live demo route
    slide = presentation.slides.add_slide(blank)
    add_header(slide, 6, "LIVE DEMO")
    add_title(
        slide,
        "ONE REAL INPUT → ONE VISIBLE RESULT",
        "Colab route: setup → load saved pipeline → predict → show validation error.",
    )
    add_box(slide, 0.65, 2.22, 5.15, 3.75, fill=LIGHT, line=BLUE)
    add_text(slide, "RAW INPUT", 0.98, 2.52, 1.55, 0.32, size=16, color=BLUE, bold=True)
    input_text = "District      Chilonzor\nSize          70 m²\nRooms         3\nLevel         3 / 5\nBuilding type Resale"
    add_text(slide, input_text, 1.00, 3.00, 4.20, 2.20, size=20)
    add_text(
        slide, "→", 5.95, 3.55, 0.70, 0.52, size=38, color=BLUE, bold=True, align=PP_ALIGN.CENTER
    )
    add_box(slide, 6.75, 2.22, 5.90, 1.72, fill=WHITE, line=GREEN)
    add_text(
        slide,
        "ESTIMATED AUGUST 2026 ASKING PRICE",
        7.05,
        2.55,
        5.25,
        0.25,
        size=14,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "$97,098 USD",
        7.05,
        2.95,
        5.25,
        0.55,
        size=31,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_box(slide, 6.75, 4.24, 5.90, 1.73, fill=WHITE, line=RED)
    add_text(slide, "EDGE CASE", 7.05, 4.52, 1.45, 0.25, size=14, color=RED, bold=True)
    add_text(
        slide, "Level 9 / max 5  →  clear ValueError", 7.05, 4.98, 5.12, 0.40, size=20, bold=True
    )
    add_text(
        slide,
        "Open: README Colab badge  •  Run all  •  Do not tour source code unless asked",
        1.15,
        6.38,
        11.05,
        0.30,
        size=15,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Slide 7 — responsible use, next step, close
    slide = presentation.slides.add_slide(blank)
    add_header(slide, 7, "CLOSE + QUESTION")
    add_title(
        slide,
        "HONEST LIMITS — CLEAR NEXT STEP",
        "A useful educational model is stronger when its boundary is explicit.",
    )
    add_card(
        slide,
        0.55,
        2.24,
        3.85,
        2.70,
        "LIMITATIONS",
        "Dated asking-price snapshot\nNo condition / exact building / legal status\nUnequal district reliability\nSource noise and market drift",
        accent=RED,
        body_size=17,
    )
    add_card(
        slide,
        4.74,
        2.24,
        3.85,
        2.70,
        "SAFE USE",
        "August 2026 reference only\nCompare with recent listings\nRequire human review\nNo lending, tax, or legal appraisal",
        accent=YELLOW,
        body_size=17,
    )
    add_card(
        slide,
        8.93,
        2.24,
        3.85,
        2.70,
        "NEXT IMPROVEMENT",
        "Collect verified transactions with exact neighborhood, condition, building year, renovation, and legal status; then use a later time holdout.",
        accent=GREEN,
        body_size=17,
    )
    add_box(slide, 2.05, 5.45, 9.20, 0.92, fill=LIGHT, line=BLUE)
    add_text(
        slide,
        "Thank you. I am ready for one evidence-based question.",
        2.35,
        5.67,
        8.60,
        0.45,
        size=18,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Slide 8 — appendix evidence map
    slide = presentation.slides.add_slide(blank)
    add_header(slide, 8, "APPENDIX / SHOW ME WHERE")
    add_title(
        slide,
        "OFFICIAL CRITERIA → EXACT REPOSITORY PROOF",
        "Open docs/capstone_evidence_matrix.md for the complete claim → location → why-proof mapping.",
    )
    evidence = [
        ("1  Problem", "README + Project Brief"),
        ("2  Data", "data/README + EDA + src/data.py"),
        ("3  Models", "model_comparison.csv + src/model.py"),
        ("4  Evaluation", "test metrics + errors + district slices"),
        ("5  Delivery", "demo.ipynb + saved joblib + predict.py"),
        ("6  Reproduce", "README + clean_run_check + CI"),
        ("7  Responsible", "README limitations + data card"),
        ("8  Defense", "deck + pitch + question bank"),
    ]
    for index, (heading, body) in enumerate(evidence):
        row, col = divmod(index, 4)
        add_card(
            slide,
            0.45 + col * 3.18,
            2.25 + row * 2.05,
            2.85,
            1.58,
            heading,
            body,
            heading_size=15,
            body_size=13,
            accent=GREEN if index not in (0, 1, 5, 7) else YELLOW,
        )
    add_text(
        slide,
        "YELLOW only where external approval/rehearsal evidence is still required.",
        2.0,
        6.52,
        9.4,
        0.30,
        size=15,
        color=YELLOW,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Slide 9 — appendix architecture
    slide = presentation.slides.add_slide(blank)
    add_header(slide, 9, "APPENDIX / ARCHITECTURE")
    add_title(
        slide,
        "THE SAME PIPELINE TRAINS AND PREDICTS",
        "Preprocessing is serialized with the estimator, preventing notebook-only inference drift.",
    )
    architecture = [
        ("RAW INPUT", "district • rooms • size • level • max_levels • new-build/resale"),
        ("VALIDATE", "required fields • positive values • level ≤ max_levels"),
        ("FEATURE", "floor_ratio = level / max_levels"),
        ("PREPROCESS", "OneHotEncoder district • numeric passthrough / scaling"),
        ("MODEL", "selected Random Forest inside sklearn Pipeline"),
        ("OUTPUT", "USD estimate + unseen/range warnings"),
    ]
    for index, (heading, body) in enumerate(architecture):
        x = 0.70 + (index % 3) * 4.18
        y = 2.20 + (index // 3) * 2.23
        add_card(
            slide,
            x,
            y,
            3.72,
            1.65,
            heading,
            body,
            body_size=14,
            accent=BLUE if index < 5 else GREEN,
        )
    add_text(
        slide,
        "Source proof: src/data.py + src/model.py + artifacts/house_price_pipeline.joblib",
        1.7,
        6.60,
        9.9,
        0.30,
        size=15,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Slide 10 — appendix Q&A route
    slide = presentation.slides.add_slide(blank)
    add_header(slide, 10, "APPENDIX / Q&A")
    add_title(
        slide,
        "ANSWER WITHOUT GUESSING",
        "Direct answer → exact evidence → limitation or next step.",
    )
    add_card(
        slide,
        0.70,
        2.25,
        3.72,
        2.30,
        "WHY RANDOM FOREST?",
        "Lowest group-safe CV MAE ($31,298).\nOpen reports/model_comparison.csv.\nTrade-off: size and interpretability.",
        accent=BLUE,
        body_size=15,
    )
    add_card(
        slide,
        4.80,
        2.25,
        3.72,
        2.30,
        "IS IT CURRENT?",
        "Yes — dated Aug 2026 asking-price snapshot.\nOpen data/README.md.\nNot a permanently live or completed-sale model.",
        accent=RED,
        body_size=15,
    )
    add_card(
        slide,
        8.90,
        2.25,
        3.72,
        2.30,
        "BIGGEST FAILURE?",
        "$1m Shayhontohur actual vs ~$234k predicted.\nOpen largest_errors.csv.\nMissing luxury/condition signal.",
        accent=YELLOW,
        body_size=15,
    )
    add_box(slide, 1.65, 5.35, 10.05, 0.92, fill=LIGHT, line=GREEN)
    add_text(
        slide,
        "Full evidence-anchored answers: docs/defense_question_bank.md",
        1.95,
        5.64,
        9.45,
        0.30,
        size=20,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(path)
